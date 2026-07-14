"""驗證 EP03 核心教學檔案處理套件。"""

from __future__ import annotations

import importlib
import sys
import tempfile
from pathlib import Path


MODULES = {
    "docx": "python-docx",
    "openpyxl": "openpyxl",
    "pptx": "python-pptx",
    "pypdf": "pypdf",
    "fitz": "PyMuPDF",
    "reportlab": "reportlab",
    "PIL": "Pillow",
    "matplotlib": "matplotlib",
    "qrcode": "qrcode",
    "markitdown": "markitdown",
}


def run_smoke_tests() -> None:
    import fitz
    import matplotlib
    import qrcode
    from docx import Document
    from markitdown import MarkItDown
    from openpyxl import Workbook
    from PIL import Image
    from pptx import Presentation
    from pptx.util import Inches
    from pypdf import PdfReader
    from reportlab.pdfgen import canvas

    matplotlib.use("Agg")
    from matplotlib import pyplot as plt

    with tempfile.TemporaryDirectory(prefix="ep03-core-") as temp_dir:
        root = Path(temp_dir)

        docx_path = root / "sample.docx"
        document = Document()
        document.add_paragraph("EP03 DOCX")
        document.save(docx_path)

        xlsx_path = root / "sample.xlsx"
        workbook = Workbook()
        workbook.active["A1"] = "EP03 XLSX"
        workbook.save(xlsx_path)

        pptx_path = root / "sample.pptx"
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1)).text = "EP03 PPTX"
        presentation.save(pptx_path)

        pdf_path = root / "sample.pdf"
        pdf_canvas = canvas.Canvas(str(pdf_path))
        pdf_canvas.drawString(72, 720, "EP03 PDF")
        pdf_canvas.save()
        assert len(PdfReader(pdf_path).pages) == 1
        with fitz.open(pdf_path) as pdf_document:
            assert pdf_document.page_count == 1
            pdf_document[0].get_pixmap(matrix=fitz.Matrix(0.5, 0.5))

        image_path = root / "sample.png"
        Image.new("RGB", (32, 32), "white").save(image_path)
        qrcode.make("https://example.com/ep03").save(root / "qr.png")

        figure = plt.figure(figsize=(2, 1))
        plt.plot([0, 1], [0, 1])
        figure.savefig(root / "chart.png")
        plt.close(figure)

        converter = MarkItDown()
        expected_text = {
            docx_path: "EP03 DOCX",
            xlsx_path: "EP03 XLSX",
            pptx_path: "EP03 PPTX",
            pdf_path: "EP03 PDF",
        }
        for source_path, expected in expected_text.items():
            result = converter.convert(str(source_path))
            assert expected in result.text_content, f"MarkItDown failed for {source_path.suffix}"


def main() -> int:
    failures: list[tuple[str, str]] = []
    for module_name, package_name in MODULES.items():
        try:
            importlib.import_module(module_name)
        except Exception as exc:
            failures.append((package_name, f"{type(exc).__name__}: {exc}"))

    if failures:
        print(f"CORE_FAIL: {len(MODULES) - len(failures)}/{len(MODULES)}")
        for package_name, error in failures:
            print(f"- {package_name}: {error}")
        return 1

    try:
        run_smoke_tests()
    except Exception as exc:
        print(f"CORE_SMOKE_FAIL: {type(exc).__name__}: {exc}")
        return 1

    version = '.'.join(str(part) for part in sys.version_info[:3])
    print(f"CORE_OK: {len(MODULES)}/{len(MODULES)} imports and file smoke tests (Python {version})")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
